import os
import re
import tempfile
from functools import lru_cache
from PIL import Image, ImageDraw, ImageFont
from typing import Dict
import numpy as np
from moviepy.editor import *
from moviepy.video.fx.all import speedx
from moviepy.audio.io.AudioFileClip import AudioFileClip
from gtts import gTTS
from pptx import Presentation
from googletrans import Translator
import warnings
import time
from tqdm import tqdm
import psutil
import gc
import glob
import subprocess
import comtypes.client
import traceback
import argparse

warnings.filterwarnings("ignore")

def translate_text(text, target_lang='en'):
    try:
        return Translator().translate(text, dest=target_lang).text
    except:
        return text
    
# 清理系统资源
gc.collect()

# 终止可能残留的ffmpeg进程
for proc in psutil.process_iter():
    if 'ffmpeg' in proc.name().lower():
        try:
            proc.kill()
        except:
            pass

class Config:
    OUTPUT_DIR = "output"
    SLIDE_IMAGE_DIR = "slides"
    LASER_COLOR = (255, 0, 0, 255)  # 带透明度的红色激光点
    LASER_RADIUS_RATIO = 0.01       # 点半径
    FADE_DURATION = 0.3             # 淡入淡出时间(秒)
    GLOW_COLOR = (255, 100, 100, 100)  # 光晕效果

class PPTSyncedConverter:
    def __init__(self):
        self.subtitle_style = {
            'font_size_ratio': 0.035,
            'bg_color': (0, 0, 128, 180),  # 半透明深蓝色背景
            'text_color': (255, 255, 255), # 白色文字
            'position_y': 0.85,            # 字幕垂直位置（0=顶部，1=底部）
            'padding': 20,                 # 文字与背景边距
            'max_width_ratio': 0.8        # 字幕最大宽度占比
        }
        self.progress_callback = None
        self._active_resources = []
        self.prs = None
        self._lock_files = set()

    def _cleanup_temp_files(self):
        """More robust temp file cleanup"""
        max_attempts = 3
        for file_path in list(self._lock_files):
            attempts = 0
            while attempts < max_attempts:
                try:                    
                    if os.path.exists(file_path):
                        os.unlink(file_path)
                    self._lock_files.remove(file_path)
                    break
                except Exception as e:
                    attempts += 1                    
                    if attempts >= max_attempts:
                        print(f"Failed to delete temp file {file_path}: {str(e)}")
                    time.sleep(0.5)
    def _cleanup_resources(self):
        """清理所有打开的媒体资源"""        
        # 先关闭所有资源
        for res in self._active_resources[:]:
            try:
                if hasattr(res, 'close'):
                    res.close()
                self._active_resources.remove(res)
            except Exception as e:
                print(f"资源释放失败: {str(e)}")
        
        # 再清理临时文件
        self._cleanup_temp_files()
        
        # 强制终止可能残留的ffmpeg进程
        try:            
            for proc in psutil.process_iter(['pid', 'name']):
                try:
                    if 'ffmpeg' in proc.info['name'].lower():
                        proc.kill()
                except:
                    pass
        except:
            pass        
        # 强制垃圾回收        
        try:
            import gc
            gc.collect()
        except:
            pass

    def _ensure_path(self, path):
        """确保目录存在且路径有效"""        
        if not path:
            raise ValueError("路径不能为空")

        dirname = os.path.dirname(path)        
        if dirname:  # 如果不是当前目录
            os.makedirs(dirname, exist_ok=True)
        return os.path.abspath(path)

    def set_progress_callback(self, callback):
        self.progress_callback = callback

    def update_progress(self, message, progress=None):
        if self.progress_callback:
            self.progress_callback(message, progress)
        else:
            print(f"[{time.strftime('%H:%M:%S')}] {message}")

    def _get_slide_notes(self, slide):
        """安全获取幻灯片备注"""
        try:
            if not slide.has_notes_slide:
                return ""
            
            notes_slide = slide.notes_slide
            text = notes_slide.notes_text_frame.text
                
            return text
        except Exception as e:
            print(f"获取备注时出错: {str(e)}")
            return ""

    def ppt_to_images(self, pptx_path, output_dir, resolution=(1920, 1080)):
        powerpoint = None
        deck = None
        try:
            # 验证输入路径
            pptx_path = self._ensure_path(pptx_path)
            output_dir = self._ensure_path(output_dir)
            
            self.update_progress("正在初始化PPT")
            comtypes.CoInitialize()
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1  # Make PowerPoint visible for debugging
            deck = powerpoint.Presentations.Open(pptx_path)

            slide_count = deck.Slides.Count
            self.update_progress(f"开始转换PPT，共{slide_count}页...")
            
            generated_images = 0
            
            for i in range(1, slide_count + 1):
                # 确保输出路径有效
                img_path = self._ensure_path(os.path.join(output_dir, f"slide_{i:03d}.png"))
                
                self.update_progress(f"正在转换第{i}/{slide_count}页...", i/slide_count)
                deck.Slides.Item(i).Export(img_path, "PNG")
                
                # 验证图片是否生成成功
                if os.path.exists(img_path) and os.path.getsize(img_path) > 0:
                    # 处理图片
                    img = Image.open(img_path).convert("RGBA")
                    if resolution != (1920, 1080):
                        img = img.resize(resolution, Image.LANCZOS)
                    img.save(img_path)
                    generated_images += 1
                else:
                    print(f"警告: 第{i}页图片生成失败")
                    
            self.update_progress("PPT转换完成!")
            return generated_images == slide_count  # 只有所有页都成功才返回True
            
        except Exception as e:
            print(f"PPT转换错误: {str(e)}")
            return False
        finally:
            if deck:
                deck.Close()
            if powerpoint:
                powerpoint.Quit()
                comtypes.CoUninitialize()

    def is_english(self, text):
        return all(ord(c) < 128 for c in text)
    
    def calculate_font_size(self, text, resolution):        
        base_size = int(resolution[1] * self.subtitle_style['font_size_ratio'])
        max_size = int(resolution[1] * 0.05)
        min_size = int(resolution[1] * 0.025)
        
        if self.is_english(text) and len(text) > 30:
            return max(min_size, base_size - 2)
        return min(max_size, base_size)
    
    def wrap_text(self, text, font, max_width):
        max_pixel_width = max_width * self.subtitle_style['max_width_ratio']
        
        if not self.is_english(text):            
            chars = list(text)
            lines = []
            current_line = ""
            for char in chars:
                test_line = current_line + char
                if font.getlength(test_line) <= max_pixel_width:                    
                    current_line = test_line
                else:
                    lines.append(current_line)
                    current_line = char
            if current_line:                
                  lines.append(current_line)
            return '\n'.join(lines)
        
        words = text.split()
        lines = []
        current_line = words[0] if words else ""
        
        for word in words[1:]:            
            test_line = f"{current_line} {word}"
            if font.getlength(test_line) <= max_pixel_width:
                current_line = test_line
            else:
                lines.append(current_line)
                current_line = word
        if current_line:
            lines.append(current_line)
        
        final_lines = []
        for line in lines:
            if font.getlength(line) > max_pixel_width:
                split_line = ""
                for char in line:                    
                    if font.getlength(split_line + char) > max_pixel_width:
                        final_lines.append(split_line)
                        split_line = char
                    else:
                        split_line += char
                if split_line:
                    final_lines.append(split_line)
            else:
                final_lines.append(line)
        
        return '\n'.join(final_lines)
    
    def process_audio_segment(self, text_segment, lang, speed, temp_dir, index):
        # Use MP3 extension consistently
        temp_audio_path = os.path.join(temp_dir, f"audio_{index}.mp3")
        final_audio_path = os.path.join(temp_dir, f"final_audio_{index}.mp3")
        
        try:
            # 1. Generate TTS audio
            tts = gTTS(text=text_segment, lang=lang, slow=False)
            tts.save(temp_audio_path)
            
            # 2. Validate and convert the audio file
            if not os.path.exists(temp_audio_path):
                raise ValueError("TTS failed to generate audio file")
                
            # 3. Use FFmpeg to ensure proper format
            ffmpeg_cmd = [
                'ffmpeg',
                '-y',  # Overwrite without asking
                '-i', temp_audio_path,
                '-acodec', 'libmp3lame',  # Use standard MP3 codec
                '-q:a', '2',  # Good quality
                '-ar', '44100',  # Standard sample rate
                '-ac', '2',  # Stereo
                final_audio_path
            ]
            
            # Run FFmpeg with error handling
            result = subprocess.run(
                ffmpeg_cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                universal_newlines=True
            )
            
            if result.returncode != 0:
                raise RuntimeError(f"FFmpeg conversion failed: {result.stderr}")
                
            # 4. Load audio with explicit format
            audio = AudioFileClip(final_audio_path, fps=44100)
            
            # 5. Speed adjustment
            if speed != 1.0:
                audio = audio.fx(speedx, speed)
                
            return audio, audio.duration
            
        except Exception as e:
            print(f"音频处理失败: {str(e)}")
            return None, 0
        finally:
            # Clean up temporary files
            for f in [temp_audio_path]:
                if os.path.exists(f):
                    try:
                        os.unlink(f)
                    except:
                        pass
    
    def create_synced_slide(self, img_path, slide_index, lang, speed, temp_dir, index, resolution):
        """创建带字幕和激光笔动画的幻灯片视频（完整中英文支持版）"""
        try:
            if self.progress_callback:
                self.progress_callback(f"PAGE_PROGRESS|{slide_index+1}|{self.prs.slides.count}|5", None)
            # ==================== 1. 初始化 ====================
            width, height = resolution
            bg_img = Image.open(img_path).convert('RGBA')
            if bg_img.size != (width, height):
                bg_img = bg_img.resize((width, height), Image.LANCZOS)
            if self.progress_callback:
                self.progress_callback(f"PAGE_PROGRESS|{slide_index+1}|{self.prs.slides.count}|10", None)

            # ==================== 2. 处理备注文本 ====================
            slide = self.prs.slides[slide_index]
            notes_text = self._get_slide_notes(slide) or ""
            if self.progress_callback:
                self.progress_callback(f"PAGE_PROGRESS|{slide_index+1}|{self.prs.slides.count}|20", None)

            # ==================== 3. 文本翻译处理 ====================
            def translate_if_needed(text):
                """非目标语言的文本自动翻译"""
                if not text:
                    return text
                    
                # 中文模式下的英文内容翻译
                if lang == 'zh-cn' and any(ord(c) > 127 for c in text):
                    try:
                        return translate_text(text, 'zh-cn')
                    except:
                        return text
                # 英文模式下的中文内容翻译
                elif lang == 'en' and any('\u4e00' <= c <= '\u9fff' for c in text):
                    try:
                        return translate_text(text, 'en')
                    except:
                        return text
                return text

            # ==================== 4. 智能分割文本 ====================
            def split_segments(text):
                """先按中文标点分段，再翻译为英文（保持语义完整）"""
                def split_by_chinese_punctuation(text):
                    """按中文标点分割文本"""
                    segments = []
                    buffer = ""
                    
                    # 特殊处理激光指令
                    temp_text = re.sub(r'(\[cursor:[^\]]+\])', r'|||\1|||', text)
                    
                    for part in temp_text.split('|||'):
                        if part.startswith('[cursor:') and part.endswith(']'):
                            if buffer:
                                segments.append(buffer)
                                buffer = ""
                            segments.append(part)
                            continue
                            
                        # 中文标点：，。！？；\n
                        parts = re.split(r'([，。！？；\n])', part)
                        
                        for p in parts:
                            if not p:
                                continue
                                
                            if p in ['，', '。', '！', '？', '；', '\n']:
                                if buffer:
                                    segments.append(buffer + p)
                                    buffer = ""
                            else:
                                buffer += p
                    
                    if buffer:
                        segments.append(buffer)
                        
                    return [s.strip() for s in segments if s.strip()]

                # 1. 先按中文标点分段
                chinese_segments = split_by_chinese_punctuation(text)
                
                # 2. 对每个分段进行翻译（如果是英文模式）
                if not lang.startswith('zh'):
                    translated_segments = []
                    for seg in chinese_segments:
                        if seg.startswith('[cursor:') and seg.endswith(']'):
                            translated_segments.append(seg)
                        else:
                            translated_segments.append(translate_text(seg, 'en'))
                    return translated_segments
                
                return chinese_segments

            segments = split_segments(notes_text)
            if self.progress_callback:
                self.progress_callback(f"PAGE_PROGRESS|{slide_index+1}|{self.prs.slides.count}|30", None)
            # ==================== 5. 生成语音和时间轴 ====================
            segment_data = []
            current_time = 0.0

            for i, segment in enumerate(segments):
                progress = 30 + int((i+1)/len(segments)*50)
                if self.progress_callback:
                    self.progress_callback(f"PAGE_PROGRESS|{slide_index+1}|{self.prs.slides.count}|{progress}", None)
                # 处理激光指令
                if segment.startswith('[cursor:') and segment.endswith(']'):
                    segment_data.append({
                        'type': 'laser',
                        'text': segment,
                        'start_time': current_time,
                        'duration': 0.1,
                        'end_time': current_time + 0.1
                    })
                    current_time += 0.1
                    continue

                # 生成语音
                clean_text = re.sub(r'\[cursor:\s*\d+,\s*\d+\]|\[cursor:\s*off\]', '', segment).strip()
                if not clean_text:
                    duration = 1.5 if slide_index == 0 else 0.5
                    segment_data.append({
                        'type': 'empty',
                        'text': segment,
                        'clean_text': "",
                        'start_time': current_time,
                        'duration': duration,
                        'end_time': current_time + duration
                    })
                    current_time += duration
                    continue

                # 确保使用正确的语言生成语音
                audio, duration = self.process_audio_segment(clean_text, lang, speed, temp_dir, f"{index}_{i}")
                duration = max(0.5, min(15.0, duration if duration else 1.0))
                
                segment_data.append({
                    'type': 'text',
                    'text': segment,
                    'clean_text': clean_text,
                    'audio': audio,
                    'start_time': current_time,
                    'duration': duration,
                    'end_time': current_time + duration
                })
                current_time += duration

            # ==================== 6. 生成激光点数据 ====================
            all_laser_points = []
            active_lasers = []

            for seg in segment_data:
                if seg['type'] != 'laser':
                    continue
                    
                text = seg['text']
                if 'off' in text.lower():
                    # 关闭激光点
                    for laser in active_lasers:
                        laser['end'] = seg['start_time']
                        all_laser_points.append(laser)
                    active_lasers = []
                else:
                    # 解析坐标（百分比转像素）
                    coords = re.findall(r'\d+', text)
                    if len(coords) >= 2:
                        x = int(width * int(coords[0]) / 100)
                        y = int(height * int(coords[1]) / 100)
                        active_lasers.append({
                            'x': max(10, min(width-10, x)),
                            'y': max(10, min(height-10, y)),
                            'start': seg['start_time'],
                            'end': float('inf')
                        })

            # 处理未关闭的激光点
            for laser in active_lasers:
                laser['end'] = segment_data[-1]['end_time']
                all_laser_points.append(laser)

            # ==================== 7. 动态渲染 ====================
            def make_frame(t):
                frame = bg_img.copy()
                overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
                draw = ImageDraw.Draw(overlay)
                
                # 1. 获取当前应显示的分段
                current_segment = next(
                    (seg for seg in segment_data 
                    if seg['start_time'] <= t < seg['end_time']),
                    None
                )
                
                # 2. 绘制字幕（仅当前分段）
                if current_segment and current_segment.get('clean_text'):
                    font_size = int(height * 0.035)
                    try:
                        font = ImageFont.truetype(
                            "simhei.ttf" if lang.startswith('zh') else "arial.ttf",
                            font_size
                        )
                    except:
                        font = ImageFont.load_default(font_size)
                    
                    # 计算文本位置（单行居中）
                    text = current_segment['clean_text']
                    text_width = font.getlength(text)
                    text_height = font_size
                    x = (width - text_width) // 2
                    y = int(height * 0.85) - text_height
                    
                    # 绘制背景
                    draw.rectangle(
                        [(x - 20, y - 10),
                        (x + text_width + 20, y + text_height + 10)],
                        fill=(0, 0, 128, 180)
                    )
                    
                    # 绘制文本
                    draw.text(
                        (x, y),
                        text,
                        font=font,
                        fill=(255, 255, 255),
                        stroke_width=1,
                        stroke_fill=(0, 0, 0)
                    )
                
                # 3. 绘制激光点
                for point in all_laser_points:
                    if point['start'] <= t < point['end']:
                        radius = int(height * 0.015)
                        # 光晕效果
                        draw.ellipse(
                            [(point['x']-radius*2, point['y']-radius*2),
                            (point['x']+radius*2, point['y']+radius*2)],
                            fill=(255, 100, 100, 100)
                        )
                        # 核心点
                        draw.ellipse(
                            [(point['x']-radius, point['y']-radius),
                            (point['x']+radius, point['y']+radius)],
                            fill=(255, 0, 0, 255)
                        )
                
                return np.array(Image.alpha_composite(frame, overlay))[:, :, :3]

            # ==================== 8. 合成最终视频 ====================
            total_duration = segment_data[-1]['end_time']
            video_clip = VideoClip(make_frame, duration=total_duration)

            if self.progress_callback:
                self.progress_callback(f"PAGE_PROGRESS|{slide_index+1}|{self.prs.slides.count}|90", None)
            
            # 添加音频
            audio_clips = [seg['audio'] for seg in segment_data if seg.get('audio')]
            if audio_clips:
                combined_audio = concatenate_audioclips(audio_clips)
                video_clip = video_clip.set_audio(combined_audio)
            
            # 写入文件
            temp_path = os.path.join(temp_dir, f"slide_{index}_final.mp4")
            video_clip.write_videofile(
                temp_path,
                fps=15,
                codec='libx264',
                audio_codec='aac',
                threads=4,
                verbose=False
            )
            
            return VideoFileClip(temp_path)

        except Exception as e:
            print(f"幻灯片创建失败: {str(e)}")
            traceback.print_exc()
            return self._create_fallback_clip(img_path, temp_dir, index)

    def _split_text_with_laser_commands(self, text):
        """分割文本并保留激光指令位置"""
        segments = []
        current_segment = ""
        
        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('[cursor:') and line.endswith(']'):
                if current_segment:
                    segments.append(current_segment)
                    current_segment = ""
                segments.append(line)
            else:
                current_segment += line + "\n"
        
        if current_segment:
            segments.append(current_segment)
        
        return [s.strip() for s in segments if s.strip()]

    def _generate_progressive_subtitle(self, text, bg_img, max_width, max_height, current_time, segment_start, segment_duration, lang = 'zh-cn', speed = 1.0):
        try:
            # Handle case where bg_img is already an Image object
            if isinstance(bg_img, Image.Image):
                img_size = bg_img.size
            else:
                # Assume it's a file path                
                temp_img = Image.open(bg_img).convert("RGBA")
                img_size = temp_img.size
                
            overlay = Image.new("RGBA", img_size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(overlay)

            line_timings = self.calculate_line_timing(text, segment_duration, lang, speed)
            if not line_timings:
                return overlay
            
            relative_time = current_time - segment_start

            visible_lines = []
            for i, line_info in enumerate(line_timings):
                if line_info['start'] <= relative_time < line_info['end']:
                    visible_lines.append(line_info['text'])
                    break
            if not visible_lines:
                if relative_time >= line_timings[-1]['end']:
                    visible_lines.append(line_timings[-1]['text'])
                else:
                    return overlay
            
            font_size = self.calculate_font_size('\n'.join(visible_lines), (max_width, max_height))            
            try:
                if self.is_english(text):
                    font = ImageFont.truetype("arial.ttf", font_size)
                else:
                    font = ImageFont.truetype("simhei.ttf", font_size)
            except:
                font = ImageFont.load_default(size=font_size)
                print("警告：使用默认字体，可能不支持中文")
            
            line_bboxes = [font.getbbox(line) for line in visible_lines]
            line_widths = [bbox[2] - bbox[0] for bbox in line_bboxes]
            line_heights = [bbox[3] - bbox[1] for bbox in line_bboxes]
            total_height = sum(line_heights)
            max_line_width = max(line_widths) if line_widths else 0
            
            bg_width = max_line_width + 2 * self.subtitle_style['padding']
            bg_height = total_height + 2 * self.subtitle_style['padding']
            y_position = img_size[1] * self.subtitle_style['position_y']
            bg_y1 = y_position - bg_height // 2
            bg_y2 = y_position + bg_height // 2
            
            draw.rectangle(
                [(img_size[0] - bg_width) // 2, bg_y1,
                (img_size[0] + bg_width) // 2, bg_y2],
                fill=self.subtitle_style['bg_color']
            )
            
            current_y = bg_y1 + self.subtitle_style['padding']
            for i, line in enumerate(visible_lines):
                text_width = font.getlength(line)
                x = (img_size[0] - text_width) // 2
                draw.text(
                    (x, current_y),
                    line,
                    font=font,
                    fill=self.subtitle_style['text_color'],
                    stroke_width=2,
                    stroke_fill=(0, 0, 0)
                )
                current_y += line_heights[i]

            return overlay        
        except Exception as e:
            print(f"生成字幕图片失败: {str(e)}")
            # Create error image            
            error_overlay = Image.new("RGBA", (max_width, max_height), (0, 0, 0, 0))
            draw = ImageDraw.Draw(error_overlay)
            draw.rectangle([(10,10), (max_width - 10, 50)], fill=(255, 0, 0, 128))
            draw.text((15,10), "Subtitle Error", fill=(255, 255, 255, 255))
            return error_overlay

    def _create_fallback_clip(self, img_path, temp_dir, index):
        """创建后备视频片段（当主流程失败时使用）"""
        try:
            if img_path and os.path.exists(img_path):
                bg_img = Image.open(img_path).convert('RGB')
                bg_clip = ImageClip(np.array(bg_img)).set_duration(3.0)
                temp_path = os.path.join(temp_dir, f"seg_{index}_fallback.mp4")
                bg_clip.write_videofile(temp_path, fps=15, verbose=False)
                return VideoFileClip(temp_path)
            else:                # 如果没有图片路径，创建空白背景
                blank = np.full((480, 854, 3), 255, dtype=np.uint8)  # 白色背景
                return ImageClip(blank).set_duration(3.0)
        except:
            # 终极后备方案
            blank = np.full((480, 854, 3), 128, dtype=np.uint8)  # 灰色背景
            return ImageClip(blank).set_duration(3.0)

    def _split_text_segments(self, text):
        """改进的文本分割方法，按照自然段落和标点分割"""
        segments = []
        
        # 首先按换行符分割
        lines = text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 中文分割规则：按句号、分号、感叹号、问号分割
            if any('\u4e00' <= c <= '\u9fff' for c in line):  # 检测中文字符
                # 使用正则表达式按中文标点分割
                parts = re.split(r'([。；！？])', line)
                
                # 重组分割后的文本
                current_segment = ""
                for part in parts:
                    if part in ['。', '；', '！', '？']:
                        if current_segment:
                            segments.append(current_segment + part)
                            current_segment = ""
                    else:
                        current_segment += part
                
                if current_segment:
                    segments.append(current_segment)
            else:
                # 英文分割规则：按句号、分号、感叹号、问号分割
                parts = re.split(r'([.;!?])', line)
                
                # 重组分割后的文本
                current_segment = ""
                for part in parts:
                    if part in ['.', ';', '!', '?']:
                        if current_segment:
                            segments.append(current_segment + part)
                            current_segment = ""
                    else:
                        current_segment += part
                
                if current_segment:
                    segments.append(current_segment)
        
        # 进一步处理过长的段落（超过30个字符）
        final_segments = []
        for segment in segments:
            segment = segment.strip()
            if not segment:
                continue
                
            # 中文处理：按逗号分割过长的段落
            if len(segment) > 30 and any('\u4e00' <= c <= '\u9fff' for c in segment):
                comma_parts = segment.split('，')
                for i, part in enumerate(comma_parts):
                    if i < len(comma_parts)-1:
                        final_segments.append(part + '，')
                    else:
                        final_segments.append(part)
            # 英文处理：按逗号分割过长的段落
            elif len(segment) > 30:
                comma_parts = segment.split(',')
                for i, part in enumerate(comma_parts):
                    if i < len(comma_parts)-1:
                        final_segments.append(part + ',')
                    else:
                        final_segments.append(part)
            else:
                final_segments.append(segment)
        
        return [s.strip() for s in final_segments if s.strip()]

    def calculate_line_timing(self, text, total_duration, lang = 'zh-cn', speed = 1.0):        
        cursor_pattern = re.compile(r'\[cursor:\s*\d+,\s*\d+\]|\[cursor:\s*off\]')
        clean_text = re.sub(cursor_pattern, '', text).strip()
        if not clean_text.strip():
            return []
        
        natural_line = clean_text.split('\n')
        natural_line = [line.strip() for line in natural_line if line.strip()]
        if not natural_line:
            return []

        final_lines = []
        for line in natural_line:
            if len(line) > 30:
                if lang and lang.startswith('zh'):
                    parts = re.split(r'([，。！？；：])', line)
                    current_part = ""
                    for part in parts:
                        if current_part and len(current_part + part) > 30:
                            if current_part.strip():
                                final_lines.append(current_part.strip())
                            current_part = part
                        else:
                            current_part += part
                    if current_part.strip():
                        final_lines.append(current_part.strip())
                else:
                    words = line.split()
                    current_line = ""
                    for word in words:
                        if current_line and len(current_line + " " + word) > 30:
                            if current_part.strip():
                                final_lines.append(current_line.strip())
                            current_line = word
                        else:
                            current_line = current_line + " " + word if current_line else word
                    if current_line.strip():
                        final_lines.append(current_line.strip())
                
            else:
                final_lines.append(line)

        if not final_lines:
            return []
        
        line_timings = []
        current_time = 0.0

        if lang and lang.startswith('zh'):
            total_chars = sum(len(line) for line in final_lines)
            char_duration = total_duration / total_chars if total_chars > 0 else 0

            for line in final_lines:
                line_char_count = len(line)
                line_duration = line_char_count * char_duration
                line_timings.append({
                    'text': line,
                    'start': current_time,
                    'end': current_time + line_duration
                })
                current_time += line_duration
        else:       
            total_words = sum(len(line.split()) for line in final_lines)
            word_duration = total_duration / total_words if total_words > 0 else 0

            for line in final_lines:
                words_in_line = len(line.split())
                line_duration = words_in_line * word_duration
                line_timings.append({
                    'text': line,
                    'start': current_time,
                    'end': current_time + line_duration
                })
                current_time += line_duration
        return line_timings

    def adjust_laser_timing(self, points, new_duration):
        if not points or new_duration <= 0:
            return points
        original_duration = max(p['end'] for p in points) if points else 0
        if original_duration <= 0:
            return points
        
        ratio = new_duration / original_duration
        adjusted_points = []
        for p in points:
            try:
                adjusted_points.append({
                    "x": p["x"],                    
                    "y": p["y"],
                    "start": p["start"] * ratio,
                    "end": p["end"] * ratio                
                })
            except:
                continue
        return adjusted_points
        
    def parse_laser_actions_precise(self, segments_data):
        """基于语音分段精确控制激光点时间"""
        laser_points = []
        active_lasers = []
        
        for seg in segments_data:
            text = seg["text"]
            start_time = seg["start_time"]
            end_time = seg["end_time"]
            
            # 检测激光指令
            if '[cursor:' in text and ']' in text:
                if 'off' in text.lower():
                    # 关闭激光点（使用当前段的开始时间）
                    for laser in active_lasers:
                        laser['end'] = start_time
                        laser_points.append(laser)
                    active_lasers = []
                else:
                    # 解析坐标（假设分辨率已在外部处理）
                    coords = re.findall(r'\d+', text)
                    if len(coords) >= 2:
                        active_lasers.append({
                            'x': int(coords[0]),
                            'y': int(coords[1]),
                            'start': start_time,  # 使用当前段的开始时间
                            'end': float('inf')   # 默认持续到被关闭
                        })
        
        # 处理未关闭的激光点
        for laser in active_lasers:
            laser['end'] = segments_data[-1]["end_time"]  # 默认持续到最后
            laser_points.append(laser)
        
        return laser_points

    def _get_timing_position(self, char_pos, total_duration):
        """计算字符位置对应的时间点"""
        # 这里可以使用更精确的时间计算逻辑
        return min(char_pos * 0.1, total_duration)  # 简单示例

    def _get_timing_position(self, char_pos, total_duration):
        """计算字符位置对应的时间点"""
        # 这里可以使用更精确的时间计算逻辑
        return min(char_pos * 0.1, total_duration)  # 简单示例
    
    def _calculate_char_timings(self, text, total_duration, lang='zh-cn', speed=1.0):       
        """计算每个字符的时间位置"""
        char_timings = {}
        if not text.strip():
            return char_timings
        # 根据语言调整计算方式+        
        if lang and lang.startswith('zh'):
            # 中文：按字符计算            
            chars = [c for c in text if c.strip()]  # 去掉空白字符
            if not chars:
                return char_timings
                
            # 平均每个字符的时间
            char_duration = total_duration / len(chars)
            
            current_time = 0.0
            char_index = 0
            
            for i, char in enumerate(text):
                if char.strip():  # 非空白字符
                    char_timings[char_index] = current_time
                    current_time += char_duration
                    char_index += 1
                # 空白字符不计入时间，但会略微增加间隔
                elif char_index > 0:
                    current_time += char_duration * 0.1
                    
        else:
            # 英文：按单词计算
            words = text.split()
            if not words:                
                return char_timings
                
            word_duration = total_duration / len(words)
            current_time = 0.0
            char_position = 0
            
            for word in words:
                # 为单词中的每个字符分配时间
                word_char_duration = word_duration / len(word) if word else 0
                for char in word:
                    char_timings[char_position] = current_time
                    current_time += word_char_duration
                    char_position += 1

                # 单词间隔
                current_time += word_duration * 0.1
                # 处理单词后的空格
                char_position += 1
        
        return char_timings
    
    def _generate_frame_with_precise_laser(self, bg_img, t, segment_data, all_laser_points, width, height, lang = 'zh-cn', speed = 1.0):
        """根据时间t生成包含精确激光点的帧"""
        try:
            # 创建背景图像副本
            frame_img = bg_img.copy()
            overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
            draw = ImageDraw.Draw(overlay)
            
            # 找到当前时间对应的段落和字幕
            current_segment = None
            for seg in segment_data:
                seg_start = seg.get("start_time", 0)
                seg_end = seg_start + seg.get("duration", 0)
                if seg_start <= t < seg_end:
                    current_segment = seg
                    break
            
            # 绘制字幕
            if current_segment and current_segment.get("clean_text"):
                try:
                    subtitle_img = self._generate_progressive_subtitle(
                        text=current_segment["clean_text"],
                        bg_img=frame_img,
                        max_width=width,
                        max_height=height,
                        current_time = t,
                        segment_start = current_segment.get("start_time", 0),
                        segment_duration = current_segment.get("duration", 0),
                        lang = lang, 
                        speed = speed
                    )
                    
                    if subtitle_img and subtitle_img.size[0] > 0 and subtitle_img.size[1] > 0:
                        y_pos = int(height * self.subtitle_style['position_y'] - subtitle_img.height // 2)
                        y_pos = max(10, min(height - subtitle_img.height - 10, y_pos))
                        x_pos = (width - subtitle_img.width) // 2                        
                        x_pos = max(10, min(width - subtitle_img.width - 10, x_pos))
                        overlay.paste(subtitle_img, (x_pos, y_pos), subtitle_img)
                except Exception as e:
                    print(f"帧{t:.2f}s字幕生成失败: {str(e)}")
            
            # 绘制激光点（基于精确时间）            
            active_laser_points = self.get_active_laser_points(all_laser_points, t)
            if active_laser_points:
                radius = int(height * Config.LASER_RADIUS_RATIO)                                
                for point in active_laser_points:               
                    try:
                        x = max(radius, min(width - radius, int(point['x'])))
                        y = max(radius, min(height - radius, int(point['y'])))
                        
                        # 绘制光晕效果
                        glow_radius = radius + 15                        
                        draw.ellipse(
                            [
                                max(0, x - glow_radius),                                
                                max(0, y - glow_radius),                               
                                min(width - 1, x + glow_radius),                                
                                min(height - 1, y + glow_radius)
                            ],
                            fill=(255, 100, 100, 150)                        
                        )
                        # 绘制激光点核心                              
                        core_radius = max(8, radius)
                        draw.ellipse(
                            [
                                max(0, x - core_radius),
                                max(0, y - core_radius),                                
                                min(width - 1, x + core_radius),
                                min(height - 1, y + core_radius)
                            ],
                           fill=(255, 0, 0, 255)
                        )
                            
                    except Exception as e:                        
                        print(f"绘制激光点失败: {str(e)}")            
            # 合并图层
            final_img = Image.alpha_composite(frame_img, overlay)
            # 转换为numpy数+            
            img_array = np.array(final_img)
            if img_array.shape[2] == 4:
                img_array = img_array[:, :, :3]
            
            return img_array
        except Exception as e:
            print(f"生成帧失败(t={t:.2f}s): {str(e)}")
            # 返回纯背景            
            bg_array = np.array(bg_img)
            if bg_array.shape[2] == 4:
                bg_array = bg_array[:, :, :3]
            return bg_array
    
    def _create_static_video_segments(self, segment_data, bg_img, all_laser_points, temp_dir, index, width, height):        
        """静态视频片段生成方法（回退方案）"""        
        try:
            video_clips = []
            
            for seg_idx, seg in enumerate(segment_data):
                try:
                    if not seg or "duration" not in seg or seg['duration'] <= 0:
                        continue
                        
                    # 创建静态帧
                    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
                    draw = ImageDraw.Draw(overlay)
                    
                    # 添加字幕
                    if seg.get("clean_text"):
                        try:
                            subtitle_img = self._generate_safe_subtitle(
                                text=seg["clean_text"],
                                bg_img = bg_img,
                                max_width=width,
                                max_height=height
                            )
                            if subtitle_img:
                                y_pos = int(height * self.subtitle_style['position_y'] - subtitle_img.height // 2)
                                y_pos = max(10, min(height - subtitle_img.height - 10, y_pos))
                                x_pos = (width - subtitle_img.width) // 2
                                x_pos = max(10, min(width - subtitle_img.width - 10, x_pos))
                                overlay.paste(subtitle_img, (x_pos, y_pos), subtitle_img)
                        except:
                            pass
                    
                    # 合并图层
                    final_img = Image.alpha_composite(bg_img, overlay)
                    img_array = np.array(final_img)
                    if img_array.shape[2] == 4:
                        img_array = img_array[:, :, :3]
                    
                    # 创建视频片段
                    clip = ImageClip(img_array).set_duration(seg["duration"])
                    if seg.get("audio"):
                        clip = clip.set_audio(seg["audio"])
                    
                    # 保存临时文件
                    temp_path = os.path.join(temp_dir, f"static_seg_{index}_{seg_idx}.mp4")
                    clip.write_videofile(temp_path, fps=15, verbose=False)
                    
                    if os.path.exists(temp_path) and os.path.getsize(temp_path) > 1024:
                        video_clips.append(VideoFileClip(temp_path))
                    
                    clip.close()
                    
                except Exception as e:
                    print(f"静态段落{seg_idx}创建失败: {str(e)}")
                    continue
            if video_clips:
                return concatenate_videoclips(video_clips)
            else:
                return self._create_fallback_clip("", temp_dir, index)

        except Exception as e:            
            print(f"静态视频生成失败: {str(e)}")
            return self._create_fallback_clip("", temp_dir, index)
        
    def get_active_laser_points(self, laser_points, current_time):
        active_points = []
        for point in laser_points:            
            if 'start' in point and 'end' in point:
                if point['start'] <= current_time < point['end']:
                    active_points.append(point)
        return active_points
    
    def _find_word_end_time(self, text_before: str, word_timings: Dict, lang: str) -> float:
        if not text_before.strip() or not word_timings:
            return 0.0
        clean_text = text_before.strip()
        if lang and lang.startswith('zh'):
            char_count = len(clean_text)
            if char_count > 0:
                for i in range(char_count - 1, -1, -1):
                    if i in word_timings:
                        return word_timings[i]['end']
        else:
            words = clean_text.split()
            if words:
                word_count = len(words)
                if word_count > 0:
                    for i in range(word_count - 1, -1, -1):
                        if i in word_timings:
                            return word_timings[i]['end']
        return 0.0
    def split_segments_with_laser_commands(text):
        segments = []
        current_segment = ""
        
        for line in text.split('\n'):
            line = line.strip()
            if line.startswith("[cursor:") and line.endswith("]"):
                if current_segment:  # 保存前一个文本段
                    segments.append({"type": "text", "content": current_segment})
                    current_segment = ""
                segments.append({"type": "laser", "command": line})  # 独立指令段
            else:
                current_segment += line + "\n"
        
        if current_segment:  # 处理最后剩余的文本段
            segments.append({"type": "text", "content": current_segment})
        
        return segments

    def convert_ppt_to_video(self, pptx_path, output_path, lang='zh-cn', resolution=(854, 480), speed=1.0, fps=15):
        try:
            # ==================== 1. 初始化阶段 (0%-5%) ====================
            self.update_progress("正在验证输入文件...", 0)
            pptx_path = self._ensure_path(pptx_path)
            output_path = self._ensure_path(output_path)

            if not os.path.exists(pptx_path):
                raise FileNotFoundError(f"PPTX文件不存在: {pptx_path}")
            if not pptx_path.lower().endswith('.pptx'):
                raise ValueError("仅支持.pptx格式的文件")

            self.update_progress("正在加载PPTX文件...", 3)
            self.prs = Presentation(pptx_path)
            if len(self.prs.slides) == 0:
                raise ValueError("PPTX中没有幻灯片")

            # ==================== 2. 创建临时工作区 ====================
            temp_dir = tempfile.mkdtemp(prefix='ppt2video_')
            self._lock_files.add(temp_dir)  # 注册临时目录以便清理
            slide_images_dir = os.path.join(temp_dir, "slides")
            os.makedirs(slide_images_dir, exist_ok=True)

            # ==================== 3. PPT转图片 (5%-20%) ====================
            self.update_progress("正在将PPT转换为图片...", 5)
            if not self.ppt_to_images(pptx_path, slide_images_dir, resolution):
                raise RuntimeError("PPT转图片失败")

            slide_images = sorted(glob.glob(os.path.join(slide_images_dir, "slide_*.png")))
            if len(slide_images) != len(self.prs.slides):
                raise RuntimeError(f"生成的图片数量({len(slide_images)})与幻灯片数量({len(self.prs.slides)})不匹配")

            # ==================== 4. 逐页处理 (20%-80%) ====================
            video_segments = []
            total_slides = len(slide_images)
            success_count = 0

            for i, img_path in enumerate(slide_images):
                # 更新页面进度 (20% + 60% * 当前进度)
                current_progress = 20 + (i / total_slides) * 60
                # 假设 total_slides 是总页数，i 是当前页索引（从0开始）
                current_page = i + 1
                total_pages = total_slides
                progress_percent = int(current_progress)  # 如果 current_progress 是 0~100 的数
                self.update_progress(f"正在处理第 {current_page}/{total_pages} 页", current_progress)

                # 🔧 关键：新增下面这行，输出 PROGRESS 格式给 PyQt 捕获
                print(f"PROGRESS|{current_page}|{total_pages}|{progress_percent}")
                sys.stdout.flush()  # 确保立即输出到管道，QProcess 才能收到

                # 内存监控（每5页检查一次）
                if i % 5 == 0:
                    mem = psutil.virtual_memory()
                    if mem.available < 200 * 1024 * 1024:
                        self.update_progress("内存不足，正在清理...", current_progress)
                        gc.collect()
                        time.sleep(1)

                try:
                    # 处理单页幻灯片（内部会调用PAGE_PROGRESS更新页内进度）
                    video_clip = self.create_synced_slide(
                        img_path=img_path,
                        slide_index=i,
                        lang=lang,
                        speed=speed,
                        temp_dir=temp_dir,
                        index=i,
                        resolution=resolution
                    )

                    # 保存临时视频片段
                    temp_video = os.path.join(temp_dir, f"clip_{i}.mp4")
                    video_clip.write_videofile(
                        temp_video,
                        fps=fps,
                        threads=2,  # 限制线程数以节省内存
                        verbose=False
                    )
                    
                    if os.path.exists(temp_video) and os.path.getsize(temp_video) > 1024:
                        video_segments.append(temp_video)
                        success_count += 1
                        self._lock_files.add(temp_video)  # 注册临时文件
                    else:
                        raise RuntimeError("生成的视频片段无效")

                except Exception as e:
                    self.update_progress(f"第 {i+1} 页处理失败: {str(e)}", current_progress)
                    continue

            # ==================== 5. 合并视频 (80%-95%) ====================
            if success_count == 0:
                raise RuntimeError("所有页面处理均失败")

            self.update_progress("正在合并视频片段...", 80)
            
            # 分段加载视频片段以避免内存不足
            final_clip = None
            try:
                for i, seg_path in enumerate(video_segments):
                    # 更新合并进度 (80% + 15% * 当前进度)
                    merge_progress = 80 + (i / len(video_segments)) * 15
                    self.update_progress(f"正在合并第 {i+1}/{len(video_segments)} 段", merge_progress)

                    clip = VideoFileClip(seg_path)
                    if final_clip is None:
                        final_clip = clip
                    else:
                        final_clip = concatenate_videoclips([final_clip, clip])

                # ==================== 6. 最终编码 (95%-100%) ====================
                self.update_progress("正在编码最终视频...", 95)
                
                # 编码进度回调函数
                def encoding_callback(progress):
                    """将moviepy的0-1进度映射到95%-100%"""
                    self.update_progress(f"编码进度: {progress*100:.1f}%", 95 + progress*5)

                print(f"[INFO] 正在导出最终视频到路径: {output_path}")
                final_clip.write_videofile(
                    output_path,
                    fps=fps,
                    threads=4,
                    preset='fast',
                    audio_codec='aac',
                    verbose=False
                )
                # 确保文件系统有足够时间写入（特别是大文件）
                time.sleep(1)

                if os.path.exists(output_path):
                    file_size = os.path.getsize(output_path)
                    if file_size > 0:
                        print(f"[SUCCESS] 视频已成功保存到: {output_path} (大小: {file_size} 字节)")
                    else:
                        print(f"[ERROR] 视频文件已创建但为空！路径: {output_path} (大小: 0 字节)")
                else:
                    print(f"[ERROR] 视频导出失败！文件未生成: {output_path}")

                self.update_progress("视频导出完成!", 100)
                return True

            finally:
                if final_clip:
                    final_clip.close()
                    
        except Exception as e:
            self.update_progress(f"转换失败: {str(e)}", -1)
            traceback.print_exc()
            return False
            
        finally:
            # 确保资源清理
            self._cleanup_resources()
            gc.collect()

if __name__ == "__main__":

    converter = PPTSyncedConverter()
    converter.subtitle_style['bg_color'] = (0, 0, 128, 180)  # 设置半透明深蓝色背景

    # 1. 解析命令行参数
    parser = argparse.ArgumentParser(description="PPT转视频工具")
    parser.add_argument('--language', type=str, default='zh-cn', help='语音语言（zh-cn/en）')
    parser.add_argument('--width', type=int, required=True, help='视频宽度（像素）')
    parser.add_argument('--height', type=int, required=True, help='视频高度（像素）')
    parser.add_argument('--speed', type=float, default=1.0, help='语音速度（0.5-2.0）')
    parser.add_argument('--input', type=str, required=True, help='输入PPTX文件路径')
    parser.add_argument('--output', type=str, default='output.mp4', help='输出视频路径')
    # 可选：如果您已经有进度格式参数
    parser.add_argument("--progress_format", type=str, default="percentage", help="进度格式")

    # 新增的参数
    parser.add_argument("--single_slide", type=int, help="仅处理指定的幻灯片编号（从1开始）")
    parser.add_argument("--slide_number", type=int, help="当前幻灯片编号（从1开始）")

    # 解析参数
    args = parser.parse_args()

    # 2. 初始化转换器
    converter = PPTSyncedConverter()
    
    # 3. 执行转换
    # 在函数调用前打印（可选，也可不加）
    print(f"[INFO] 接收到的视频输出路径为: {args.output}")
    success = converter.convert_ppt_to_video(
        pptx_path=args.input,
        output_path=args.output,
        lang=args.language,
        resolution=(args.width, args.height),
        speed=args.speed
    )
    
    if success:
        print("转换成功！")
    else:
        print("转换失败")
