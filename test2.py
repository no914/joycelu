import os
import re
import tempfile
from functools import lru_cache
from PIL import Image, ImageDraw, ImageFont
from typing import List, Dict
import numpy as np
from moviepy.editor import *
from moviepy.video.fx.all import speedx
from moviepy.audio.io.AudioFileClip import AudioFileClip
from gtts import gTTS
from pptx import Presentation
import comtypes.client
from googletrans import Translator
import warnings
import time
from tqdm import tqdm
import psutil
import gc
import glob
import subprocess
import cv2

warnings.filterwarnings("ignore")

def translate_text(text, target_lang='en'):
    try:
        return Translator().translate(text, dest=target_lang).text
    except:
        return text
    
# 清理系统资源
gc.collect()

# 终止可能残留的ffmpeg进程
for proc in psutil.process_iter():
    if 'ffmpeg' in proc.name().lower():
        try:
            proc.kill()
        except:
            pass

class Config:
    OUTPUT_DIR = "output"
    SLIDE_IMAGE_DIR = "slides"
    LASER_COLOR = (255, 0, 0, 255)  # 带透明度的红色激光点
    LASER_RADIUS_RATIO = 0.04       # 点半径
    FADE_DURATION = 0.3             # 淡入淡出时间(秒)
    GLOW_COLOR = (255, 100, 100, 100)  # 光晕效果

class PPTSyncedConverter:
    def __init__(self):
        self.subtitle_style = {
            'font_size_ratio': 0.035,
            'bg_color': (0, 0, 128, 180),  # 半透明深蓝色背景
            'text_color': (255, 255, 255), # 白色文字
            'position_y': 0.85,            # 字幕垂直位置（0=顶部，1=底部）
            'padding': 20,                 # 文字与背景边距
            'max_width_ratio': 0.8        # 字幕最大宽度占比
        }
        self.progress_callback = None
        self._active_resources = []
        self.prs = None
        self._lock_files = set()

    def _cleanup_temp_files(self):
        """More robust temp file cleanup"""
        max_attempts = 3
        for file_path in list(self._lock_files):
            attempts = 0
            while attempts < max_attempts:
                try:
                    if os.path.exists(file_path):
                        os.unlink(file_path)
                    self._lock_files.remove(file_path)
                    break
                except Exception as e:
                    attempts += 1
                    if attempts >= max_attempts:
                        print(f"Failed to delete temp file {file_path}: {str(e)}")
                    time.sleep(0.5)

    def _cleanup_resources(self):
        """清理所有打开的媒体资源"""
        # 先关闭所有资源
        for res in self._active_resources[:]:
            try:
                if hasattr(res, 'close'):
                    res.close()
                self._active_resources.remove(res)
            except Exception as e:
                print(f"资源释放失败: {str(e)}")
        
        # 再清理临时文件
        self._cleanup_temp_files()
        
        # 强制终止可能残留的ffmpeg进程
        try:
            for proc in psutil.process_iter(['pid', 'name']):
                try:
                    if 'ffmpeg' in proc.info['name'].lower():
                        proc.kill()
                except:
                    pass
        except:
            pass
        
        # 强制垃圾回收
        try:
            import gc
            gc.collect()
        except:
            pass

    def _ensure_path(self, path):
        """确保目录存在且路径有效"""
        if not path:
            raise ValueError("路径不能为空")
        
        dirname = os.path.dirname(path)
        if dirname:  # 如果不是当前目录
            os.makedirs(dirname, exist_ok=True)
        return os.path.abspath(path)

    def set_progress_callback(self, callback):
        self.progress_callback = callback

    def update_progress(self, message, progress=None):
        if self.progress_callback:
            self.progress_callback(message, progress)
        else:
            print(f"[{time.strftime('%H:%M:%S')}] {message}")

    def _get_slide_notes(self, slide):
        """安全获取幻灯片备注"""
        try:
            if not slide.has_notes_slide:
                return ""
            
            notes_slide = slide.notes_slide
            text = notes_slide.notes_text_frame.text
            print(f"✅ 提取到的备注内容: '{text}'")  # 调试输出
                
            return text
        except Exception as e:
            print(f"获取备注时出错: {str(e)}")
            return ""

    def ppt_to_images(self, pptx_path, output_dir, resolution=(1920, 1080)):
        powerpoint = None
        try:
            # 验证输入路径
            pptx_path = self._ensure_path(pptx_path)
            output_dir = self._ensure_path(output_dir)
            
            self.update_progress("正在初始化PPT")
            comtypes.CoInitialize()
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            deck = powerpoint.Presentations.Open(pptx_path)

            slide_count = deck.Slides.Count
            self.update_progress(f"开始转换PPT，共{slide_count}页...")
            
            for i in range(1, slide_count + 1):
                # 确保输出路径有效
                img_path = self._ensure_path(os.path.join(output_dir, f"slide_{i:03d}.png"))
                
                self.update_progress(f"正在转换第{i}/{slide_count}页...", i/slide_count)
                deck.Slides.Item(i).Export(img_path, "PNG")
                
                # 处理图片
                img = Image.open(img_path).convert("RGBA")
                if resolution != (1920, 1080):
                    img = img.resize(resolution, Image.LANCZOS)
                img.save(img_path)
            
            self.update_progress("PPT转换完成!")
            return True
        except Exception as e:
            print(f"PPT转换错误: {str(e)}")
            return False
        finally:
            if powerpoint:
                deck.Close()
                powerpoint.Quit()
                comtypes.CoUninitialize()

    def is_english(self, text):
        return all(ord(c) < 128 for c in text)
    
    def calculate_font_size(self, text, resolution):
        base_size = int(resolution[1] * self.subtitle_style['font_size_ratio'])
        max_size = int(resolution[1] * 0.05)
        min_size = int(resolution[1] * 0.025)
        
        if self.is_english(text) and len(text) > 30:
            return max(min_size, base_size - 2)
        return min(max_size, base_size)
    
    def wrap_text(self, text, font, max_width):
        max_pixel_width = max_width * self.subtitle_style['max_width_ratio']
        
        if not self.is_english(text):
            chars = list(text)
            lines = []
            current_line = ""
            for char in chars:
                test_line = current_line + char
                if font.getlength(test_line) <= max_pixel_width:
                    current_line = test_line
                else:
                    lines.append(current_line)
                    current_line = char
            if current_line:
                lines.append(current_line)
            return '\n'.join(lines)
        
        words = text.split()
        lines = []
        current_line = words[0] if words else ""
        
        for word in words[1:]:
            test_line = f"{current_line} {word}"
            if font.getlength(test_line) <= max_pixel_width:
                current_line = test_line
            else:
                lines.append(current_line)
                current_line = word
        if current_line:
            lines.append(current_line)
        
        final_lines = []
        for line in lines:
            if font.getlength(line) > max_pixel_width:
                split_line = ""
                for char in line:
                    if font.getlength(split_line + char) > max_pixel_width:
                        final_lines.append(split_line)
                        split_line = char
                    else:
                        split_line += char
                if split_line:
                    final_lines.append(split_line)
            else:
                final_lines.append(line)
        
        return '\n'.join(final_lines)
    
    def process_audio_segment(self, text_segment, lang, speed, temp_dir, index):
        # Use MP3 extension consistently
        try:
            if not text_segment or not text_segment.strip():
                duration = 1.0
                return None, duration
            
            clean_text = text_segment.strip()
            speed = max(0.5, min(3.0, float(speed))) if speed else 1.0

            estimated_duration = max(1.0, len(clean_text) * 0.15)
            estimated_duration = min(15.0, estimated_duration)

            print(f"尝试生成音频：文本长度={len(clean_text)}，预估时长={estimated_duration:.1f}秒")
            temp_audio_path = os.path.join(temp_dir, f"audio_{index}.mp3")

            # 1. Generate TTS audio
            try:
                tts = gTTS(text=clean_text, lang=lang, slow=False)
                tts.save(temp_audio_path)

                if not os.path.exists(temp_audio_path) or os.path.getsize(temp_audio_path) < 1024:
                    raise ValueError("TTS生成的文件无效或太小")
                
                try:
                    audio = AudioFileClip(temp_audio_path)
                    if audio.duration <= 0 or audio.duration > 30:
                        audio.close()
                        raise ValueError(f"音频时常异常: {audio.duration}")

                    # 5. Speed adjustment
                    if speed != 1.0:
                        try:
                            audio = audio.fx(speedx, speed)
                        except Exception as e:
                            print(f"语速调整失败，使用原速: {str(e)}")

                    final_duration = audio.duration
                    if final_duration <= 0 or final_duration > 30:
                        audio.close()
                        raise ValueError(f"最终音频时长异常：{final_duration}")
                    print(f"音频生成成功：时长={final_duration:.1f}秒")
                    return audio, final_duration

                except Exception as e:
                    print(f"音频加载失败: {str(e)}")
                    return None, estimated_duration
            except Exception as e:
                print(f"TTS生成失败：{str(e)}")
                return None, estimated_duration
        except Exception as e:
            print(f"音频处理完全失败: {str(0)}")
            fallback_duration = 2.0
            return None, fallback_duration
        
        finally:
            # Clean up temporary files
            temp_files = [
                os.path.join(temp_dir, f"audio_{index}.mp3"),
                os.path.join(temp_dir, f"final_audio_{index}.mp3")
            ]
            for f in temp_files:
                if os.path.exists(f):
                    try:
                        os.unlink(f)
                    except:
                        pass
    
    def create_synced_slide(self, img_path, slide_index, lang, speed, temp_dir, index, resolution):
        """创建带字幕和激光笔动画的幻灯片视频（完整安全版本）"""
        try:
            # ==================== 1. 初始化检查 ====================
            gc.collect()
            if psutil.virtual_memory().percent > 85:
                time.sleep(5)

            # 验证输入参数
            if not os.path.exists(img_path):
                raise FileNotFoundError(f"幻灯片图片不存在: {img_path}")
            
            width, height = resolution
            if width <= 0 or height <= 0:
                raise ValueError(f"无效的分辨率: {width}x{height}")

            # ==================== 2. 加载幻灯片内容 ====================
            slide = self.prs.slides[slide_index]
            notes_text = self._get_slide_notes(slide) or ""
            
            # 提取纯净文本（移除激光指令）
            display_text = re.sub(r'\[cursor:\s*\d+,\s*\d+\]|\[cursor:\s*off\]', '', notes_text).strip()

            # ==================== 3. 文本预处理 ====================
            # 自动翻译非目标语言文本
            if display_text and lang != 'zh-cn' and any('\u4e00' <= c <= '\u9fff' for c in display_text):
                display_text = translate_text(display_text, 'en')
            
            # 分割文本为适合语音合成的段落
            segments = self._split_text_segments(notes_text)
            if not segments:
                segments = [display_text] if display_text else [""]

            # ==================== 4. 生成语音和时长数据 ====================
            video_clips = []
            segment_data = []
            current_time = 0.0

            for i, segment in enumerate(segments):
                try:
                    clean_segment = re.sub(r'\[cursor:\s*\d+,\s*\d+\]|\[cursor:\s*off\]', '', segment).strip()

                    audio = None
                    duration = 2.0
                    
                    # 生成语音（带错误处理）
                    if clean_segment:
                        try:
                            audio, duration = self.process_audio_segment(
                                clean_segment, lang, speed, temp_dir, f"{index}_{i}"
                            )
                            print(f"段落{i}: 音频={'有效' if audio else '无'}，时长={duration:.1f}秒")
                        except Exception as e:
                            print(f"段落{i} 音频生成异常：{str(e)}")
                            audio = None
                            duration = max(1.0, len(clean_segment) * 0.15)
                    else:
                        duration = 3.0 if slide_index == 0 else 1.5  # 默认时长
                        print(f"段落{i}: 空文本，使用默认时长{duration:.1f}秒")

                    duration = max(0.5, min(15.0, duration))

                    segment_data.append({
                        "text": segment,
                        "clean_text": clean_segment,
                        "audio": audio,
                        "duration": duration,
                        "start_time": current_time,
                        "end_time": current_time + duration
                    })
                    current_time += duration

                except Exception as e:
                    print(f"段落{i}处理失败: {str(e)}")
                    fallback_duration = 2.0
                    segment_data.append({
                        "text": segment if 'segment' in locals() else "fallback",
                        "clean_text": "",
                        "audio": None,
                        "duration": fallback_duration,
                        "start_time": current_time,
                        "end_time": current_time + fallback_duration
                    })
                    current_time += fallback_duration

            # ==================== 5. 生成视频片段和精确激光点时间 ====================
            all_laser_points = []
            print(f"原始备注文本: {notes_text}")
            print(f"分割后的段落: {segments}")
            
            for segment in segment_data:
                # 使用新的精确时间解析方法
                segment_laser_points = self.parse_laser_actions_precise(
                    text=segment["text"],
                    segment_start=segment["start_time"],
                    actual_audio_duration=segment["duration"],
                    lang=lang,
                    speed=speed,
                    slide_width=width,
                    slide_height=height
                )
                print(f"基于段落解析的激光点: {segment_laser_points}")
                all_laser_points.extend(segment_laser_points)

            try:
                bg_img = Image.open(img_path).convert('RGBA')
                if bg_img.size != (width, height):
                    bg_img = bg_img.resize((width, height), Image.LANCZOS)
            except Exception as e:
                print(f"背景图片处理失败: {str(e)}")
                bg_img = Image.new('RGBA', (width, height), (255, 255, 255, 255))

                            # ==================== 6. 创建动态激光点视频 ====================
                # 使用动态帧生成来实现精确的激光点时间控制
                total_duration = sum(seg.get("duration", 0) for seg in segment_data)
                
                if total_duration > 0:
                    try:
                        # 创建动态视频剪辑
                        def make_frame(t):
                            return self._generate_frame_with_precise_laser(
                                bg_img, t, segment_data, all_laser_points, width, height
                            )
                        
                        # 创建视频剪辑
                        video_clip = VideoClip(make_frame, duration=total_duration)
                        
                        # 添加音频
                        audio_clips = []
                        for seg in segment_data:
                            if seg.get("audio") and hasattr(seg["audio"], 'duration'):
                                try:
                                    audio_clips.append(seg["audio"])
                                except Exception as e:
                                    print(f"音频处理失败: {str(e)}")
                        
                        if audio_clips:
                            try:
                                combined_audio = concatenate_audioclips(audio_clips)
                                video_clip = video_clip.set_audio(combined_audio)
                            except Exception as e:
                                print(f"音频合并失败: {str(e)}")
                        
                        return video_clip
                        
                    except Exception as e:
                        print(f"动态视频生成失败: {str(e)}")
                        # 回退到静态方法
                        return self._create_static_video_segments(segment_data, bg_img, all_laser_points, temp_dir, index, width, height)
                
                # 如果没有有效的段落数据，返回空
                return self._create_fallback_clip(img_path, temp_dir, index)

        except Exception as e:
            print(f"创建幻灯片时出错: {str(e)}")
            return self._create_fallback_clip(img_path, temp_dir, index)

    def _generate_safe_subtitle(self, text, bg_img, max_width, max_height):
        try:
            # Handle case where bg_img is already an Image object
            if isinstance(bg_img, Image.Image):
                img_size = bg_img.size
            else:
                # Assume it's a file path
                temp_img = Image.open(bg_img).convert("RGBA")
                img_size = temp_img.size
                
            overlay = Image.new("RGBA", img_size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(overlay)
            
            font_size = self.calculate_font_size(text, (max_width, max_height))
            
            try:
                if self.is_english(text):
                    font = ImageFont.truetype("arial.ttf", font_size)
                else:
                    font = ImageFont.truetype("simhei.ttf", font_size)
            except:
                font = ImageFont.load_default(size=font_size)
                print("警告：使用默认字体，可能不支持中文")

            wrapped_text = self.wrap_text(text, font, max_width)
            lines = wrapped_text.split('\n')
            
            line_bboxes = [font.getbbox(line) for line in lines]
            line_widths = [bbox[2] - bbox[0] for bbox in line_bboxes]
            line_heights = [bbox[3] - bbox[1] for bbox in line_bboxes]
            total_height = sum(line_heights)
            max_line_width = max(line_widths)
            
            bg_width = max_line_width + 2 * self.subtitle_style['padding']
            bg_height = total_height + 2 * self.subtitle_style['padding']
            y_position = img_size[1] * self.subtitle_style['position_y']
            bg_y1 = y_position - bg_height // 2
            bg_y2 = y_position + bg_height // 2
            
            draw.rectangle(
                [(img_size[0] - bg_width) // 2, bg_y1,
                (img_size[0] + bg_width) // 2, bg_y2],
                fill=self.subtitle_style['bg_color']
            )
            
            current_y = bg_y1 + self.subtitle_style['padding']
            for i, line in enumerate(lines):
                text_width = font.getlength(line)
                x = (img_size[0] - text_width) // 2
                draw.text(
                    (x, current_y),
                    line,
                    font=font,
                    fill=self.subtitle_style['text_color'],
                    stroke_width=2,
                    stroke_fill=(0, 0, 0)
                )
                current_y += line_heights[i]

            return overlay
        
        except Exception as e:
            print(f"生成字幕图片失败: {str(e)}")
            # Create error image
            error_overlay = Image.new("RGBA", (max_width, max_height), (0, 0, 0, 0))
            draw = ImageDraw.Draw(error_overlay)
            draw.rectangle([(10,10), (max_width - 10, 50)], fill=(255, 0, 0, 128))
            draw.text((15,10), "Subtitle Error", fill=(255, 255, 255, 255))
            return error_overlay

    def _create_fallback_clip(self, img_path, temp_dir, index):
        """创建后备视频片段（当主流程失败时使用）"""
        try:
            if img_path and os.path.exists(img_path):
                bg_img = Image.open(img_path).convert('RGB')
                bg_clip = ImageClip(np.array(bg_img)).set_duration(3.0)
                temp_path = os.path.join(temp_dir, f"seg_{index}_fallback.mp4")
                bg_clip.write_videofile(temp_path, fps=15, verbose=False)
                return VideoFileClip(temp_path)
            else:
                # 如果没有图片路径，创建空白背景
                blank = np.full((480, 854, 3), 255, dtype=np.uint8)  # 白色背景
                return ImageClip(blank).set_duration(3.0)
        except:
            # 终极后备方案
            blank = np.full((480, 854, 3), 128, dtype=np.uint8)  # 灰色背景
            return ImageClip(blank).set_duration(3.0)

    def _split_text_segments(self, text):
        """辅助方法：分割文本为语音段落"""
        segments = []
        current_segment = ""
        cursor_pattern = re.compile(r'(\[cursor:\s*\d+,\s*\d+\]|\[cursor:\s*off\])')
        
        for line in text.split('\n'):
            if line.strip():
                parts = cursor_pattern.split(line)
                for part in parts:
                    if not part:
                        continue
                        
                    if cursor_pattern.match(part):
                        current_segment += part
                        if current_segment.strip():
                            segments.append(current_segment)
                            current_segment = ""
                    else:
                        current_segment += part
                        if part.endswith(('。', '!', '?', ';')):
                            segments.append(current_segment)
                            current_segment = ""
        
        if current_segment:
            segments.append(current_segment)
        
        return segments

    def calculate_word_timing(self, text: str, lang: str = 'zh-cn', speed: float = 1.0) -> Dict:
        clean_text = re.sub(r'\[cursor:\s*\d+,\s*\d+\]|\[cursor:\s*off\]', '', text).strip()
        if not clean_text:
            return {}
        speed = max(0.5, min(3.0, float(speed))) if speed else 1.0
        if lang and lang.startswith('zh'):
            chars = list(clean_text)
            chars_per_second = 3.5 * speed
            
            word_timings = {}
            current_time = 0.0

            for i, char in enumerate(chars):
                if char.strip():
                    char_duration = 1.0
                    word_timings[i] = {
                        'word': char,
                        'start': current_time,
                        'end': current_time + char_duration
                    }
                    current_time += char_duration
                else:
                    current_time += 0.1
        else:
            words = clean_text.split()
            words_per_second = 2.5 * speed

            word_timings = {}
            current_time = 0.0

            for i, word in enumerate(words):
                word_duration = len(word) / (words_per_second * 3)
                word_duration = max(0.2, min(1,0, word_duration))

                word_timings[i] = {
                    'word': word,
                    'start': current_time,
                    'end': current_time + word_duration
                }
                current_time += word_duration + 0.1
        return word_timings

    def adjust_laser_timing(self, points, new_duration):
        if not points or new_duration <= 0:
            return points
        
        original_duration = max(p['end'] for p in points) if points else 0
        if original_duration <= 0:
            return points
        
        ratio = new_duration / original_duration
        adjusted_points = []
        for p in points:
            try:
                adjusted_points.append({
                    "x": p["x"],
                    "y": p["y"],
                    "start": p["start"] * ratio,
                    "end": p["end"] * ratio
                })
            except:
                continue
        return adjusted_points
        
    def parse_laser_actions_precise(self, text, segment_start, actual_audio_duration, lang='zh-cn', speed=1.0, slide_width=1920, slide_height=1080):
        """精确的激光点时间解析 - 基于实际语音时长"""
        cursor_positions = []
        
        # 正则表达式匹配激光点指令
        cursor_pattern = re.compile(
            r'\[cursor:\s*(\d+)\s*,\s*(\d+)\]|\[cursor:\s*off\]',
            re.IGNORECASE
        )
        
        # 移除激光点指令，获取纯文本用于时间计算
        clean_text = re.sub(cursor_pattern, '', text)
        
        if not clean_text.strip():
            return []
        
        # 计算每个字符的时间位置（基于实际音频时长）
        char_timings = self._calculate_char_timings(clean_text, actual_audio_duration, lang, speed)
        
        # 解析激光点指令
        for match in cursor_pattern.finditer(text):
            cursor_cmd = match.group()
            
            # 计算指令前的文本长度（不包括之前的指令）
            text_before_cursor = text[:match.start()]
            clean_text_before = re.sub(cursor_pattern, '', text_before_cursor)
            char_count_before = len(clean_text_before.strip())
            
            # 计算激光点出现的精确时间
            if char_count_before > 0:
                # 基于前面文本的字符数量计算时间
                appear_time = segment_start + char_timings.get(char_count_before - 1, actual_audio_duration)
            else:
                appear_time = segment_start
            
            print(f"激光点指令: {cursor_cmd}")
            print(f"指令前文本: '{clean_text_before}' (长度: {char_count_before})")
            print(f"计算出现时间: {appear_time:.2f}s (段落开始: {segment_start:.2f}s)")
            
            if 'off' in cursor_cmd.lower():
                # 关闭激光点
                if cursor_positions:
                    cursor_positions[-1]['end'] = appear_time
                    print(f"激光点关闭时间: {appear_time:.2f}s")
            else:
                # 解析坐标
                coords = [int(g) for g in match.groups() if g]
                if len(coords) >= 2:
                    x = int(slide_width * coords[0] / 100)
                    y = int(slide_height * coords[1] / 100)
                    
                    # 添加新激光点
                    laser_point = {
                        'x': x,
                        'y': y,
                        'start': appear_time,
                        'end': segment_start + actual_audio_duration  # 默认持续到段落结束
                    }
                    cursor_positions.append(laser_point)
                    print(f"新激光点: 坐标({x}, {y}), 开始时间: {appear_time:.2f}s")
        
        return cursor_positions
    
    def _calculate_char_timings(self, text, total_duration, lang='zh-cn', speed=1.0):
        """计算每个字符的时间位置"""
        char_timings = {}
        
        if not text.strip():
            return char_timings
        
        # 根据语言调整计算方式
        if lang and lang.startswith('zh'):
            # 中文：按字符计算
            chars = [c for c in text if c.strip()]  # 去掉空白字符
            if not chars:
                return char_timings
                
            # 平均每个字符的时间
            char_duration = total_duration / len(chars)
            
            current_time = 0.0
            char_index = 0
            
            for i, char in enumerate(text):
                if char.strip():  # 非空白字符
                    char_timings[char_index] = current_time
                    current_time += char_duration
                    char_index += 1
                # 空白字符不计入时间，但会略微增加间隔
                elif char_index > 0:
                    current_time += char_duration * 0.1
                    
        else:
            # 英文：按单词计算
            words = text.split()
            if not words:
                return char_timings
                
            word_duration = total_duration / len(words)
            current_time = 0.0
            char_position = 0
            
            for word in words:
                # 为单词中的每个字符分配时间
                word_char_duration = word_duration / len(word) if word else 0
                for char in word:
                    char_timings[char_position] = current_time
                    current_time += word_char_duration
                    char_position += 1
                
                # 单词间隔
                current_time += word_duration * 0.1
                # 处理单词后的空格
                char_position += 1
        
        return char_timings
    
    def _generate_frame_with_precise_laser(self, bg_img, t, segment_data, all_laser_points, width, height):
        """根据时间t生成包含精确激光点的帧"""
        try:
            # 创建背景图像副本
            frame_img = bg_img.copy()
            overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
            draw = ImageDraw.Draw(overlay)
            
            # 找到当前时间对应的段落和字幕
            current_segment = None
            for seg in segment_data:
                seg_start = seg.get("start_time", 0)
                seg_end = seg_start + seg.get("duration", 0)
                if seg_start <= t < seg_end:
                    current_segment = seg
                    break
            
            # 绘制字幕
            if current_segment and current_segment.get("clean_text"):
                try:
                    subtitle_img = self._generate_safe_subtitle(
                        text=current_segment["clean_text"],
                        bg_img=frame_img,
                        max_width=width,
                        max_height=height
                    )
                    
                    if subtitle_img and subtitle_img.size[0] > 0 and subtitle_img.size[1] > 0:
                        y_pos = int(height * self.subtitle_style['position_y'] - subtitle_img.height // 2)
                        y_pos = max(10, min(height - subtitle_img.height - 10, y_pos))
                        x_pos = (width - subtitle_img.width) // 2
                        x_pos = max(10, min(width - subtitle_img.width - 10, x_pos))
                        overlay.paste(subtitle_img, (x_pos, y_pos), subtitle_img)
                except Exception as e:
                    print(f"帧{t:.2f}s字幕生成失败: {str(e)}")
            
            # 绘制激光点（基于精确时间）
            active_laser_points = self.get_active_laser_points(all_laser_points, t)
            if active_laser_points:
                radius = int(height * Config.LASER_RADIUS_RATIO)
                
                for point in active_laser_points:
                    try:
                        x = max(radius, min(width - radius, int(point['x'])))
                        y = max(radius, min(height - radius, int(point['y'])))
                        
                        # 绘制光晕效果
                        glow_radius = radius + 15
                        draw.ellipse(
                            [
                                max(0, x - glow_radius),
                                max(0, y - glow_radius),
                                min(width - 1, x + glow_radius),
                                min(height - 1, y + glow_radius)
                            ],
                            fill=(255, 100, 100, 150)
                        )
                        
                        # 绘制激光点核心
                        core_radius = max(8, radius)
                        draw.ellipse(
                            [
                                max(0, x - core_radius),
                                max(0, y - core_radius),
                                min(width - 1, x + core_radius),
                                min(height - 1, y + core_radius)
                            ],
                            fill=(255, 0, 0, 255)
                        )
                        
                        # 添加白色中心点
                        center_radius = max(3, core_radius // 3)
                        draw.ellipse(
                            [
                                x - center_radius,
                                y - center_radius,
                                x + center_radius,
                                y + center_radius
                            ],
                            fill=(255, 255, 255, 255)
                        )
                        
                        if t % 2 < 0.1:  # 每2秒打印一次，避免日志过多
                            print(f"时间{t:.2f}s: 激光点活跃 坐标({x}, {y})")
                            
                    except Exception as e:
                        print(f"绘制激光点失败: {str(e)}")
            
            # 合并图层
            final_img = Image.alpha_composite(frame_img, overlay)
            
            # 转换为numpy数组
            img_array = np.array(final_img)
            if img_array.shape[2] == 4:  # RGBA转RGB
                img_array = img_array[:, :, :3]
            
            return img_array
            
        except Exception as e:
            print(f"生成帧失败(t={t:.2f}s): {str(e)}")
            # 返回纯背景
            bg_array = np.array(bg_img)
            if bg_array.shape[2] == 4:
                bg_array = bg_array[:, :, :3]
            return bg_array
    
    def _create_static_video_segments(self, segment_data, bg_img, all_laser_points, temp_dir, index, width, height):
        """静态视频片段生成方法（回退方案）"""
        try:
            video_clips = []
            
            for seg_idx, seg in enumerate(segment_data):
                try:
                    if not seg or "duration" not in seg or seg['duration'] <= 0:
                        continue
                        
                    # 创建静态帧
                    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
                    draw = ImageDraw.Draw(overlay)
                    
                    # 添加字幕
                    if seg.get("clean_text"):
                        try:
                            subtitle_img = self._generate_safe_subtitle(
                                text=seg["clean_text"],
                                bg_img=bg_img,
                                max_width=width,
                                max_height=height
                            )
                            if subtitle_img:
                                y_pos = int(height * self.subtitle_style['position_y'] - subtitle_img.height // 2)
                                y_pos = max(10, min(height - subtitle_img.height - 10, y_pos))
                                x_pos = (width - subtitle_img.width) // 2
                                x_pos = max(10, min(width - subtitle_img.width - 10, x_pos))
                                overlay.paste(subtitle_img, (x_pos, y_pos), subtitle_img)
                        except:
                            pass
                    
                    # 合并图层
                    final_img = Image.alpha_composite(bg_img, overlay)
                    img_array = np.array(final_img)
                    if img_array.shape[2] == 4:
                        img_array = img_array[:, :, :3]
                    
                    # 创建视频片段
                    clip = ImageClip(img_array).set_duration(seg["duration"])
                    if seg.get("audio"):
                        clip = clip.set_audio(seg["audio"])
                    
                    # 保存临时文件
                    temp_path = os.path.join(temp_dir, f"static_seg_{index}_{seg_idx}.mp4")
                    clip.write_videofile(temp_path, fps=15, verbose=False)
                    
                    if os.path.exists(temp_path) and os.path.getsize(temp_path) > 1024:
                        video_clips.append(VideoFileClip(temp_path))
                    
                    clip.close()
                    
                except Exception as e:
                    print(f"静态段落{seg_idx}创建失败: {str(e)}")
                    continue
            
            if video_clips:
                return concatenate_videoclips(video_clips)
            else:
                return self._create_fallback_clip("", temp_dir, index)
                
        except Exception as e:
            print(f"静态视频生成失败: {str(e)}")
            return self._create_fallback_clip("", temp_dir, index)
        
    def get_active_laser_points(self, laser_points, current_time):
        active_points = []
        for point in laser_points:
            if point['start'] <= current_time <= point['end']:
                active_points.append(point)
        return active_points
        
    def _find_word_end_time(self, text_before: str, word_timings: Dict, lang: str) -> float:
        if not text_before.strip() or not word_timings:
            return 0.0
        clean_text = text_before.strip()

        if lang and lang.startswith('zh'):
            char_count = len(clean_text)
            if char_count > 0:
                for i in range(char_count - 1, -1, -1):
                    if i in word_timings:
                        return word_timings[i]['end']
        else:
            words = clean_text.split()
            if words:
                word_count = len(words)
                if word_count > 0:
                    for i in range(word_count - 1, -1, -1):
                        if i in word_timings:
                            return word_timings[i]['end']
        return 0.0

    def convert_ppt_to_video(self, pptx_path, output_path, lang='zh-cn', resolution=(854, 480), speed=1.0, fps=15):
        try:
            # ==================== 1. 初始化验证 ====================
            self.update_progress("正在验证输入文件...")
            pptx_path = self._ensure_path(pptx_path)
            output_path = self._ensure_path(output_path)
            
            if not os.path.exists(pptx_path):
                raise FileNotFoundError(f"PPTX文件不存在: {pptx_path}")
                
            if not pptx_path.lower().endswith('.pptx'):
                raise ValueError("仅支持.pptx格式的文件")

            # ==================== 2. 加载PPTX ====================
            self.update_progress("正在加载PPTX文件...")
            try:
                self.prs = Presentation(pptx_path)
                if len(self.prs.slides) == 0:
                    raise ValueError("PPTX中没有幻灯片")
            except Exception as e:
                raise ValueError(f"PPTX加载失败: {str(e)}")

            # ==================== 3. 创建临时工作区 ====================
            temp_dir = tempfile.mkdtemp(prefix='ppt2video_')
            try:
                temp_dir = self._ensure_path(temp_dir)
                self.update_progress(f"临时工作区: {temp_dir}")
                
                # ==================== 4. PPT转图片 ====================
                slide_images = []
                self.update_progress("正在将PPT转换为图片...")
                if not self.ppt_to_images(pptx_path, temp_dir, resolution):
                    raise RuntimeError("PPT转图片失败")
                
                slide_images = sorted(glob.glob(os.path.join(temp_dir, "slide_*.png")))
                if len(slide_images) != len(self.prs.slides):
                    raise RuntimeError(f"生成的图片数量({len(slide_images)})与幻灯片数量({len(self.prs.slides)})不匹配")

                # ==================== 5. 逐页处理 ====================
                video_segments = []
                success_count = 0
                
                for i, img_path in enumerate(slide_images):
                    try:
                        # 内存监控
                        mem = psutil.virtual_memory()
                        if mem.available < 200 * 1024 * 1024:
                            self.update_progress("内存不足，暂停处理...")
                            time.sleep(5)
                            gc.collect()
                        
                        # 处理单页
                        self.update_progress(f"正在处理第{i+1}/{len(slide_images)}页...", (i+1)/len(slide_images))
                        
                        result = self.create_synced_slide(
                            img_path=img_path,
                            slide_index=i,
                            lang=lang,
                            speed=speed,
                            temp_dir=temp_dir,
                            index=i,
                            resolution=resolution
                        )
                        
                        if result:
                            temp_video = os.path.join(temp_dir, f"clip_{i}.mp4")
                            result.write_videofile(temp_video, fps=fps)
                            if os.path.exists(temp_video) and os.path.getsize(temp_video) > 0:
                                video_segments.append(temp_video)
                                success_count += 1
                            else:
                                self.update_progress(f"生成的第{i+1}页视频无效")
                        else:
                            self.update_progress(f"跳过无效的第{i+1}页")
                            
                    except Exception as e:
                        self.update_progress(f"第{i+1}页处理失败: {str(e)}")
                        continue

                # ==================== 6. 合并视频 ====================
                if success_count > 0:
                    self.update_progress("正在合并视频片段...")
                    
                    # 使用moviepy合并视频
                    clips = []
                    for seg in video_segments:
                        try:
                            clip = VideoFileClip(seg)
                            clips.append(clip)
                        except Exception as e:
                            self.update_progress(f"加载视频片段失败: {seg} - {str(e)}")
                    
                    if clips:
                        final_clip = concatenate_videoclips(clips)
                        final_clip.write_videofile(
                            output_path,
                            fps=fps,
                            threads=4,
                            preset='ultrafast',
                            audio_codec='aac',
                            verbose=False
                        )
                        final_clip.close()
                        
                        for clip in clips:
                            clip.close()
                        
                        return True
                    
                return False
                
            finally:
                # 清理临时目录
                try:
                    for root, dirs, files in os.walk(temp_dir, topdown=False):
                        for name in files:
                            try:
                                os.unlink(os.path.join(root, name))
                            except:
                                pass
                        for name in dirs:
                            try:
                                os.rmdir(os.path.join(root, name))
                            except:
                                pass
                    os.rmdir(temp_dir)
                except:
                    pass
                
        except Exception as e:
            self.update_progress(f"转换失败: {str(e)}")
            return False
        finally:
            self._cleanup_resources()

def get_users_choice():
    print("请选择语音语言")
    print("1. 中文")
    print("2. 英文")
    while True:
        lang_choice = input("请输入选项(1-2): ").strip()
        if lang_choice in ('1', '2'):
            lang = 'zh-cn' if lang_choice == '1' else 'en'
            break
        print("无效输入，请重新选择")

    print("\n请选择分辨率")
    print("1. 1920x1080 (全高清)")
    print("2. 1280x720 (高清)")
    while True:
        resolution_choice = input("请输入选项(1-2): ").strip()
        if resolution_choice == '1':
            resolution = (1920, 1080)
            break
        elif resolution_choice == '2':
            resolution = (1280, 720)
            break
        else:
            print("无效输入，请重新选择")
        
    print("\n请选择语音速度")
    print("1. 慢速 (0.75x)")
    print("2. 正常 (1x)")
    print("3. 快速 (1.5x)")
    while True:
        speed_choice = input("请输入选项(1-3): ").strip()
        if speed_choice == '1':
            speed = 0.75
            break
        elif speed_choice == '2':
            speed = 1
            break
        elif speed_choice == '3':
            speed = 1.5
            break
        else:
            print("无效输入，请重新选择")

    return lang, resolution, speed

if __name__ == "__main__":

    converter = PPTSyncedConverter()
    converter.subtitle_style['bg_color'] = (0, 0, 128, 180)  # 设置半透明深蓝色背景

    lang, resolution, speed = get_users_choice()

    def progress_callback(message, progress=None):
        if progress is not None:
            print(f"[{time.strftime('%H:%M:%S')}] {message} ({progress*100:.1f}%)")
        else:
            print(f"[{time.strftime('%H:%M:%S')}] {message}")
    
    converter.set_progress_callback(progress_callback)
    
    success = converter.convert_ppt_to_video(
        pptx_path="test.pptx",
        output_path="output.mp4",
        lang=lang,
        resolution=resolution, 
        speed=speed
    )
    print("转换成功!" if success else "转换失败")
